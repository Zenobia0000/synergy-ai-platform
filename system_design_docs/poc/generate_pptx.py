"""
Synergy AI POC Pitch Deck Generator
Produces a Google-style clean presentation (15 main slides + 4 appendix)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors (Google-inspired clean palette) ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
DARK = RGBColor(0x20, 0x21, 0x24)
GRAY_600 = RGBColor(0x5F, 0x63, 0x68)
GRAY_400 = RGBColor(0x9A, 0xA0, 0xA6)
BLUE = RGBColor(0x42, 0x85, 0xF4)       # Google Blue — primary accent
GREEN = RGBColor(0x34, 0xA8, 0x53)       # success / check
RED = RGBColor(0xEA, 0x43, 0x35)         # alert
YELLOW = RGBColor(0xFB, 0xBC, 0x05)      # warning
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFE)
ACCENT_BG = RGBColor(0xF1, 0xF3, 0xF4)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ──

def _add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_shape_bg(slide, left, top, width, height, color, corner_radius=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape

def _tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def _set_text(tf, text, size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Noto Sans TC"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return p

def _add_para(tf, text, size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Noto Sans TC"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p

def _add_bullet(tf, text, size=16, color=DARK, level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Noto Sans TC"
    p.font.bold = bold
    p.level = level
    p.space_before = Pt(4)
    return p

def _add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    return table

def _set_cell(table, row, col, text, size=13, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, bg=None):
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Noto Sans TC"
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

def _header_row(table, texts, size=13):
    for i, t in enumerate(texts):
        _set_cell(table, 0, i, t, size=size, bold=True, color=WHITE, bg=BLUE)

def _page_number(slide, num, total=19):
    tb = _tb(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4))
    _set_text(tb.text_frame, f"{num} / {total}", size=10, color=GRAY_400, alignment=PP_ALIGN.RIGHT)

def _section_tag(slide, text):
    """Small colored tag in top-left"""
    shape = _add_shape_bg(slide, Inches(0.6), Inches(0.4), Inches(1.8), Inches(0.35), BLUE, corner_radius=0.15)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Noto Sans TC"
    p.alignment = PP_ALIGN.CENTER

# ══════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_add_bg(s, WHITE)

# Blue accent bar at top
_add_shape_bg(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)

# Main title
tb = _tb(s, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2))
_set_text(tb.text_frame, "Synergy AI", size=54, bold=True, color=DARK)

tb = _tb(s, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8))
_set_text(tb.text_frame, "健康管理事業 AI 營運平台", size=28, color=GRAY_600)

tb = _tb(s, Inches(1.5), Inches(4.0), Inches(10), Inches(0.6))
_set_text(tb.text_frame, "POC 提案簡報", size=20, color=BLUE, bold=True)

# Tagline
_add_shape_bg(s, Inches(1.5), Inches(5.2), Inches(5.5), Inches(0.6), LIGHT_BLUE_BG, corner_radius=0.1)
tb = _tb(s, Inches(1.7), Inches(5.25), Inches(5), Inches(0.5))
_set_text(tb.text_frame, "先吸引、再分流、後複製", size=18, color=BLUE, bold=True)

tb = _tb(s, Inches(1.5), Inches(6.4), Inches(5), Inches(0.4))
_set_text(tb.text_frame, "2026-04-22", size=14, color=GRAY_400)

# ══════════════════════════════════════════════════
# SLIDE 2 — Pain Points
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "痛點")
_page_number(s, 2)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "你的痛點，我們都聽見了", size=36, bold=True, color=DARK)

# Three cards
cards = [
    ("📝", "行銷內容", "手動寫文、品質不穩、產量有限", "曝光不夠，錯過熱點時機"),
    ("🔍", "潛客篩選", "憑經驗判斷、容易誤判", "高潛力客戶流失，時間浪費"),
    ("📊", "團隊管理", "口頭追蹤、沒有數據", "成功模式無法複製"),
]
for i, (icon, title, problem, loss) in enumerate(cards):
    x = Inches(0.8 + i * 4.0)
    card = _add_shape_bg(s, x, Inches(2.2), Inches(3.6), Inches(3.8), NEAR_WHITE, corner_radius=0.03)
    # icon
    tb = _tb(s, x + Inches(0.3), Inches(2.5), Inches(3), Inches(0.6))
    _set_text(tb.text_frame, icon, size=36)
    # title
    tb = _tb(s, x + Inches(0.3), Inches(3.1), Inches(3), Inches(0.5))
    _set_text(tb.text_frame, title, size=22, bold=True, color=DARK)
    # problem
    tb = _tb(s, x + Inches(0.3), Inches(3.7), Inches(3), Inches(0.8))
    _set_text(tb.text_frame, problem, size=14, color=GRAY_600)
    # loss
    tb = _tb(s, x + Inches(0.3), Inches(4.7), Inches(3), Inches(0.8))
    _set_text(tb.text_frame, f"→ {loss}", size=13, color=RED, bold=True)

tb = _tb(s, Inches(0.8), Inches(6.4), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "核心矛盾：事業要規模化，但每件事都還在靠「人力撐」。", size=16, bold=True, color=GRAY_600)

# ══════════════════════════════════════════════════
# SLIDE 3 — Before / After
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "問題 → 解法")
_page_number(s, 3)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "AI 導入前後對比", size=36, bold=True, color=DARK)

table = _add_table(s, 6, 3, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.0))
_header_row(table, ["情境", "現況", "AI 導入後"])
data = [
    ("一篇社群貼文", "手動 30–60 分鐘", "AI 產出 < 30 秒"),
    ("問卷評估", "人工判讀 15–20 分鐘", "AI 即時分析 < 15 秒"),
    ("名單分級", "靠感覺、靠經驗", "量化評分，準確率 ≥ 85%"),
    ("目標追蹤", "口頭確認 / Excel", "即時儀表板，隨時掌握"),
    ("新人上手", "數週摸索", "系統引導，30 分鐘上手"),
]
for r, (a, b, c) in enumerate(data, 1):
    bg = NEAR_WHITE if r % 2 == 0 else WHITE
    _set_cell(table, r, 0, a, bold=True, bg=bg)
    _set_cell(table, r, 1, b, color=RED, bg=bg)
    _set_cell(table, r, 2, c, color=GREEN, bold=True, bg=bg)

tb = _tb(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "這不是未來式，這是 13 週後的樣子。", size=18, bold=True, color=BLUE)

# ══════════════════════════════════════════════════
# SLIDE 4 — Three Module Framework
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "解法框架")
_page_number(s, 4)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "三大 AI 模組 — 先吸引、再分流、後複製", size=32, bold=True, color=DARK)

# Three connected boxes
modules_data = [
    (BLUE, "模組 1", "AI 獲客行銷", "AI 生成貼文 → 多平台自動發布", "先吸引"),
    (GREEN, "模組 2", "AI 問卷評估", "隱私問卷 → AI 分析 → 名單分級", "再分流"),
    (RGBColor(0x7B, 0x1F, 0xA2), "模組 6", "AI 團隊管理", "目標追蹤 → 績效看板 → 培訓", "後複製"),
]
for i, (accent, num, name, desc, tagline) in enumerate(modules_data):
    y = Inches(2.2 + i * 1.65)
    # Accent bar
    _add_shape_bg(s, Inches(1.0), y, Inches(0.12), Inches(1.35), accent)
    # Content
    tb = _tb(s, Inches(1.4), y + Inches(0.05), Inches(5), Inches(0.4))
    _set_text(tb.text_frame, f"{num}：{name}", size=22, bold=True, color=DARK)
    tb = _tb(s, Inches(1.4), y + Inches(0.5), Inches(5), Inches(0.4))
    _set_text(tb.text_frame, desc, size=15, color=GRAY_600)
    # Tag
    tag = _add_shape_bg(s, Inches(8.5), y + Inches(0.2), Inches(2.0), Inches(0.5), accent, corner_radius=0.15)
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Noto Sans TC"
    p.alignment = PP_ALIGN.CENTER

    # Arrow between boxes
    if i < 2:
        tb = _tb(s, Inches(5.5), y + Inches(1.2), Inches(1), Inches(0.5))
        _set_text(tb.text_frame, "↓", size=24, color=GRAY_400, alignment=PP_ALIGN.CENTER)

tb = _tb(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "三個模組各自獨立可用，整合後形成完整業務飛輪。", size=16, color=GRAY_600)

# ══════════════════════════════════════════════════
# SLIDE 5 — Module 1 Detail
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "模組 1")
_page_number(s, 5)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "AI 獲客行銷", size=36, bold=True, color=DARK)
tb = _tb(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.4))
_set_text(tb.text_frame, "讓內容產出從「耗時」變「秒出」", size=18, color=BLUE)

tb = _tb(s, Inches(0.8), Inches(2.4), Inches(6), Inches(3.5))
tf = tb.text_frame
_set_text(tf, "功能亮點", size=20, bold=True, color=DARK)
for txt in [
    "AI 自動生成貼文文案 + 配圖建議",
    "多平台排程發布：FB / IG / LINE 一鍵搞定",
    "SEO 落地頁自動生成，吸引自然流量",
    "內容行事曆：活動預排不遺漏",
]:
    _add_bullet(tf, f"•  {txt}", size=15, color=GRAY_600)

# Metrics card
_add_shape_bg(s, Inches(7.5), Inches(2.4), Inches(5), Inches(3.0), LIGHT_BLUE_BG, corner_radius=0.03)
tb = _tb(s, Inches(7.8), Inches(2.6), Inches(4.5), Inches(0.4))
_set_text(tb.text_frame, "衡量指標", size=18, bold=True, color=BLUE)
metrics = [("產出速度", "提升 ≥ 50%"), ("月內容產量", "3 倍（不增人力）"), ("發布一致性", "達成率 ≥ 90%")]
for j, (k, v) in enumerate(metrics):
    tb = _tb(s, Inches(7.8), Inches(3.2 + j * 0.65), Inches(4.5), Inches(0.5))
    tf = tb.text_frame
    _set_text(tf, k, size=14, bold=True, color=DARK)
    _add_para(tf, v, size=16, bold=True, color=BLUE, space_before=Pt(2))

# ══════════════════════════════════════════════════
# SLIDE 6 — Module 2 Detail
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "模組 2")
_page_number(s, 6)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "AI 問卷評估", size=36, bold=True, color=DARK)
tb = _tb(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.4))
_set_text(tb.text_frame, "讓每一個潛客都得到精準的第一印象", size=18, color=GREEN)

tb = _tb(s, Inches(0.8), Inches(2.4), Inches(6), Inches(3.5))
tf = tb.text_frame
_set_text(tf, "功能亮點", size=20, bold=True, color=DARK)
for txt in [
    "隱私友善的分段式問卷（不嚇跑使用者）",
    "AI 即時健康評估 + 個人化產品推薦",
    "名單自動分級：高 / 中 / 低潛力",
    "智慧分流規則引擎：依評分自動觸發行動",
]:
    _add_bullet(tf, f"•  {txt}", size=15, color=GRAY_600)

_add_shape_bg(s, Inches(7.5), Inches(2.4), Inches(5), Inches(3.0), RGBColor(0xE6, 0xF4, 0xEA), corner_radius=0.03)
tb = _tb(s, Inches(7.8), Inches(2.6), Inches(4.5), Inches(0.4))
_set_text(tb.text_frame, "衡量指標", size=18, bold=True, color=GREEN)
metrics = [("問卷完成率", "≥ 70%"), ("AI 推薦準確率", "≥ 85%"), ("轉化率", "較現況提升")]
for j, (k, v) in enumerate(metrics):
    tb = _tb(s, Inches(7.8), Inches(3.2 + j * 0.65), Inches(4.5), Inches(0.5))
    tf = tb.text_frame
    _set_text(tf, k, size=14, bold=True, color=DARK)
    _add_para(tf, v, size=16, bold=True, color=GREEN, space_before=Pt(2))

# Badge
badge = _add_shape_bg(s, Inches(7.8), Inches(5.2), Inches(4.2), Inches(0.45), GREEN, corner_radius=0.15)
tf = badge.text_frame
p = tf.paragraphs[0]
p.text = "✅ POC v1 已完成 — 核心引擎已驗證"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Noto Sans TC"
p.alignment = PP_ALIGN.CENTER

# ══════════════════════════════════════════════════
# SLIDE 7 — Module 6 Detail
# ══════════════════════════════════════════════════
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "模組 6")
_page_number(s, 7)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "AI 團隊管理", size=36, bold=True, color=DARK)
tb = _tb(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.4))
_set_text(tb.text_frame, "讓「成功模式」可以被看見、被複製", size=18, color=PURPLE)

tb = _tb(s, Inches(0.8), Inches(2.4), Inches(6), Inches(3.5))
tf = tb.text_frame
_set_text(tf, "功能亮點", size=20, bold=True, color=DARK)
for txt in [
    "個人 + 團隊月目標設定",
    "即時績效儀表板（不用等月底對帳）",
    "MEGA 培訓進度追蹤",
    "落後自動告警（提前預警）",
]:
    _add_bullet(tf, f"•  {txt}", size=15, color=GRAY_600)

_add_shape_bg(s, Inches(7.5), Inches(2.4), Inches(5), Inches(3.0), RGBColor(0xF3, 0xE8, 0xFD), corner_radius=0.03)
tb = _tb(s, Inches(7.8), Inches(2.6), Inches(4.5), Inches(0.4))
_set_text(tb.text_frame, "衡量指標", size=18, bold=True, color=PURPLE)
metrics = [("主管週使用率", "≥ 70%"), ("目標達成率", "提升對比"), ("新人完訓時間", "縮短 ≥ 30%")]
for j, (k, v) in enumerate(metrics):
    tb = _tb(s, Inches(7.8), Inches(3.2 + j * 0.65), Inches(4.5), Inches(0.5))
    tf = tb.text_frame
    _set_text(tf, k, size=14, bold=True, color=DARK)
    _add_para(tf, v, size=16, bold=True, color=PURPLE, space_before=Pt(2))

# ══════════════════════════════════════════════════
# SLIDE 8 — Already Built (Trust Anchor)
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "信任基礎")
_page_number(s, 8)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "不是從零開始 — 我們已經走了這麼遠", size=32, bold=True, color=DARK)

table = _add_table(s, 4, 3, Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.8))
_header_row(table, ["項目", "已驗證 ✅", "POC 擴展方向 →"])
rows = [
    ("模組 2 POC v1", "問卷 → AI 分析 → 產品推薦", "+ 名單分級 + 分流引擎"),
    ("模組 1 基礎架構", "多平台分發 + n8n 自動化", "+ AI 內容生成 + 排程"),
    ("技術棧", "FastAPI + React + LiteLLM + Gemini + PostgreSQL", "+ 模組 6 全新開發"),
]
for r, (a, b, c) in enumerate(rows, 1):
    bg = NEAR_WHITE if r % 2 == 0 else WHITE
    _set_cell(table, r, 0, a, bold=True, bg=bg)
    _set_cell(table, r, 1, b, color=GREEN, bg=bg)
    _set_cell(table, r, 2, c, color=BLUE, bold=True, bg=bg)

# Key message box
_add_shape_bg(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.3), LIGHT_BLUE_BG, corner_radius=0.03)
tb = _tb(s, Inches(1.2), Inches(5.5), Inches(10.5), Inches(1.0))
tf = tb.text_frame
_set_text(tf, "這代表什麼？", size=18, bold=True, color=BLUE)
for txt in [
    "✅  技術可行性已驗證 — 不是概念，是跑得動的系統",
    "✅  風險在已知範圍內 — 不是在黑盒子裡摸索",
    "✅  POC 的每一週都在真實交付物上推進",
]:
    _add_bullet(tf, txt, size=14, color=DARK)

# ══════════════════════════════════════════════════
# SLIDE 9 — Timeline
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "時程")
_page_number(s, 9)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "13 週 · 5 階段 · 5 個里程碑", size=36, bold=True, color=DARK)

# Timeline bar
phases = [
    ("W1", "基礎建設", BLUE, 1),
    ("W2–W4", "模組 2 補強", GREEN, 3),
    ("W5–W7", "模組 1 補強", BLUE, 3),
    ("W8–W11", "模組 6 開發", PURPLE, 4),
    ("W12–13", "整合+UAT", YELLOW, 2),
]
x_start = Inches(0.8)
bar_y = Inches(2.6)
unit_w = Inches(0.88)  # per week
for weeks_label, name, color, weeks in phases:
    w = unit_w * weeks
    bar = _add_shape_bg(s, x_start, bar_y, w, Inches(0.7), color, corner_radius=0.05)
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Noto Sans TC"
    p.alignment = PP_ALIGN.CENTER
    # Week label below
    tb = _tb(s, x_start, bar_y + Inches(0.75), w, Inches(0.3))
    _set_text(tb.text_frame, weeks_label, size=10, color=GRAY_400, alignment=PP_ALIGN.CENTER)
    x_start += w + Inches(0.08)

# Milestone table
table = _add_table(s, 6, 3, Inches(0.8), Inches(3.8), Inches(11.5), Inches(2.8))
_header_row(table, ["里程碑", "時間", "交付物"])
ms = [
    ("M1", "W1 (05/04)", "共通基礎就緒（CI/CD、DB、認證）"),
    ("M2", "W4 (05/25)", "模組 2 完整走通 — 業務可實際試用"),
    ("M3", "W7 (06/15)", "模組 1 完整走通 — 第一批 AI 貼文"),
    ("M4", "W11 (07/13)", "模組 6 完整走通 — 儀表板上線"),
    ("M5", "W13 (07/27)", "三模組整合 · UAT · POC Demo"),
]
for r, (a, b, c) in enumerate(ms, 1):
    bg = NEAR_WHITE if r % 2 == 0 else WHITE
    _set_cell(table, r, 0, a, bold=True, color=BLUE, bg=bg)
    _set_cell(table, r, 1, b, bg=bg)
    _set_cell(table, r, 2, c, bg=bg)

tb = _tb(s, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4))
_set_text(tb.text_frame, "每個里程碑都是真實可展示的系統，不是 PPT。", size=14, bold=True, color=BLUE)

# ══════════════════════════════════════════════════
# SLIDE 10 — Collaboration Model
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "合作模式")
_page_number(s, 10)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "合夥共創 — 不是傳統甲乙方", size=36, bold=True, color=DARK)

table = _add_table(s, 6, 2, Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.2))
_header_row(table, ["面向", "做法"])
rows = [
    ("需求定義", "雙方共同定義，不是單方丟規格書"),
    ("驗收方式", "每階段 Gate Review，雙方共同確認"),
    ("決策機制", "業務方向：客戶拍板 ｜ 技術實作：開發拍板"),
    ("同步節奏", "每週定期 Check-in（30 分鐘）"),
    ("問題處理", "Slack / 即時通訊，24 小時內回應"),
]
for r, (a, b) in enumerate(rows, 1):
    bg = NEAR_WHITE if r % 2 == 0 else WHITE
    _set_cell(table, r, 0, a, bold=True, bg=bg)
    _set_cell(table, r, 1, b, bg=bg)

_add_shape_bg(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8), LIGHT_BLUE_BG, corner_radius=0.03)
tb = _tb(s, Inches(1.2), Inches(5.9), Inches(10.5), Inches(0.6))
_set_text(tb.text_frame, "你的業務知識 ＋ 我們的技術能力 ＝ 真正可用的系統", size=18, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 11 — Pricing
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "投資")
_page_number(s, 11)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "投資與報價", size=36, bold=True, color=DARK)

# Left: Fixed fee
_add_shape_bg(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), NEAR_WHITE, corner_radius=0.03)
tb = _tb(s, Inches(1.1), Inches(2.2), Inches(5), Inches(0.5))
_set_text(tb.text_frame, "方案 A — 固定費用", size=20, bold=True, color=DARK)

items_a = [
    ("POC 開發費", "NT$ _______"),
    ("技術基礎設施", "~NT$ 10,000–30,000"),
    ("（LLM API + 雲端主機）", "實報實銷"),
]
for j, (k, v) in enumerate(items_a):
    tb = _tb(s, Inches(1.1), Inches(3.0 + j * 0.6), Inches(4.8), Inches(0.5))
    tf = tb.text_frame
    _set_text(tf, k, size=14, color=GRAY_600)
    run = tf.paragraphs[0].add_run()
    run.text = f"    {v}"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = DARK
    run.font.name = "Noto Sans TC"

tb = _tb(s, Inches(1.1), Inches(5.0), Inches(4.8), Inches(0.5))
_set_text(tb.text_frame, "合計：NT$ _______", size=22, bold=True, color=BLUE)

# Right: Staged payment
_add_shape_bg(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5), LIGHT_BLUE_BG, corner_radius=0.03)
badge = _add_shape_bg(s, Inches(10.3), Inches(2.1), Inches(2.0), Inches(0.35), BLUE, corner_radius=0.15)
tf = badge.text_frame
p = tf.paragraphs[0]
p.text = "推薦"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Noto Sans TC"
p.alignment = PP_ALIGN.CENTER

tb = _tb(s, Inches(7.1), Inches(2.2), Inches(5), Inches(0.5))
_set_text(tb.text_frame, "方案 B — 分階段付款", size=20, bold=True, color=DARK)

table = _add_table(s, 5, 3, Inches(7.1), Inches(3.0), Inches(5.2), Inches(2.6))
_header_row(table, ["節點", "比例", "觸發條件"], size=12)
payment = [
    ("簽約", "30%", "Charter 簽核"),
    ("M2", "25%", "模組 2 走通"),
    ("M4", "25%", "模組 6 走通"),
    ("M5", "20%", "Demo + 驗收"),
]
for r, (a, b, c) in enumerate(payment, 1):
    bg = WHITE
    _set_cell(table, r, 0, a, bold=True, size=12, bg=bg)
    _set_cell(table, r, 1, b, size=12, color=BLUE, bold=True, bg=bg, alignment=PP_ALIGN.CENTER)
    _set_cell(table, r, 2, c, size=12, bg=bg)

tb = _tb(s, Inches(7.1), Inches(5.8), Inches(5), Inches(0.5))
_set_text(tb.text_frame, "每一筆錢都對應一個真實交付物。", size=14, bold=True, color=BLUE)

# ══════════════════════════════════════════════════
# SLIDE 12 — ROI
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "ROI")
_page_number(s, 12)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "為什麼值得這筆投資", size=36, bold=True, color=DARK)

roi_cards = [
    ("💰", "人力節省", "每位經營者每天省 1–2 小時\n10 人團隊，月省 15 萬工時成本", BLUE),
    ("📈", "轉化率提升", "結構化問卷 + AI 分級\n轉化率即使只提升 5%\n業績以倍數計算", GREEN),
    ("🔄", "複製速度", "新人 30 分鐘上手\n成功模式數據化\n可複製到下一個市場", PURPLE),
    ("🗄️", "數據資產", "問卷 + 行為數據\n是公司長期資產\nPOC 數據直接延用", YELLOW),
]
for i, (icon, title, desc, accent) in enumerate(roi_cards):
    x = Inches(0.8 + i * 3.1)
    _add_shape_bg(s, x, Inches(2.2), Inches(2.8), Inches(4.2), NEAR_WHITE, corner_radius=0.03)
    _add_shape_bg(s, x, Inches(2.2), Inches(2.8), Inches(0.08), accent)
    tb = _tb(s, x + Inches(0.2), Inches(2.5), Inches(2.4), Inches(0.5))
    _set_text(tb.text_frame, f"{icon}  {title}", size=18, bold=True, color=DARK)
    tb = _tb(s, x + Inches(0.2), Inches(3.2), Inches(2.4), Inches(2.5))
    _set_text(tb.text_frame, desc, size=13, color=GRAY_600)

# ══════════════════════════════════════════════════
# SLIDE 13 — Risk Management
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "風險控管")
_page_number(s, 13)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "風險控管 — 你的錢花在刀口上", size=36, bold=True, color=DARK)

protections = [
    ("分階段付款", "每個里程碑交付並驗收後\n才觸發下一筆付款\n資金風險是分散的", BLUE),
    ("Go / No-Go 機制", "POC 結束有明確決策框架\n不會因為已付錢\n就被綁架繼續投入", GREEN),
    ("Pivot 選項", "某個方向跑不通\n我們有能力調整\n而非全部重來", YELLOW),
]
for i, (title, desc, accent) in enumerate(protections):
    x = Inches(0.8 + i * 4.1)
    _add_shape_bg(s, x, Inches(2.2), Inches(3.7), Inches(2.8), NEAR_WHITE, corner_radius=0.03)
    _add_shape_bg(s, x, Inches(2.2), Inches(3.7), Inches(0.08), accent)
    tb = _tb(s, x + Inches(0.3), Inches(2.5), Inches(3.2), Inches(0.5))
    _set_text(tb.text_frame, title, size=20, bold=True, color=DARK)
    tb = _tb(s, x + Inches(0.3), Inches(3.2), Inches(3.2), Inches(1.5))
    _set_text(tb.text_frame, desc, size=14, color=GRAY_600)

# Decision tree
_add_shape_bg(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5), ACCENT_BG, corner_radius=0.03)
tb = _tb(s, Inches(1.2), Inches(5.5), Inches(10.5), Inches(1.3))
tf = tb.text_frame
_set_text(tf, "POC Demo 後的決策路徑", size=16, bold=True, color=DARK)
_add_bullet(tf, "✅  三模組達標 → 進入正式版（加 3 個模組）", size=13, color=GREEN, bold=True)
_add_bullet(tf, "🔄  部分達標 → 4 週改善衝刺後重新評估", size=13, color=YELLOW, bold=True)
_add_bullet(tf, "⏹️  未達標 → 停止投入，交付技術報告 + 學習成果", size=13, color=RED, bold=True)

# ══════════════════════════════════════════════════
# SLIDE 14 — Client Responsibilities
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_section_tag(s, "共同責任")
_page_number(s, 14)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "你需要準備什麼", size=36, bold=True, color=DARK)

phases_resp = [
    ("Phase 0\n4/28 前", "領域知識文件\n核心業務指標定義\n指定 1 位業務窗口", BLUE),
    ("Phase 1\nW2–W4", "問卷題目確認\n分流規則定義", GREEN),
    ("Phase 2\nW5–W7", "社群內容風格指引\n各平台帳號權限", BLUE),
    ("Phase 3\nW8–W11", "培訓計畫素材\n營運指標定義", PURPLE),
]
for i, (phase, items, accent) in enumerate(phases_resp):
    x = Inches(0.8 + i * 3.1)
    _add_shape_bg(s, x, Inches(2.2), Inches(2.8), Inches(3.5), NEAR_WHITE, corner_radius=0.03)
    _add_shape_bg(s, x, Inches(2.2), Inches(2.8), Inches(0.08), accent)
    tb = _tb(s, x + Inches(0.2), Inches(2.5), Inches(2.4), Inches(0.7))
    _set_text(tb.text_frame, phase, size=14, bold=True, color=accent)
    tb = _tb(s, x + Inches(0.2), Inches(3.4), Inches(2.4), Inches(2.0))
    _set_text(tb.text_frame, items, size=14, color=GRAY_600)

_add_shape_bg(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.7), LIGHT_BLUE_BG, corner_radius=0.03)
tb = _tb(s, Inches(1.2), Inches(6.15), Inches(10.5), Inches(0.5))
_set_text(tb.text_frame, "「不是甩手給我們做，是一起把它做好。」", size=18, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 15 — CTA / Closing
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, WHITE)
_add_shape_bg(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
_page_number(s, 15)

tb = _tb(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "下一步 — 一起開始", size=40, bold=True, color=DARK)

table = _add_table(s, 5, 3, Inches(1.5), Inches(2.5), Inches(10), Inches(2.8))
_header_row(table, ["步驟", "行動", "時間"])
steps = [
    ("1", "確認 POC 範圍 → 簽核 Charter", "本週內"),
    ("2", "確認報價與付款方式", "本週內"),
    ("3", "正式 Kick-off Meeting", "2026-04-28"),
    ("4", "第一個里程碑：基礎就緒", "2026-05-04"),
]
for r, (a, b, c) in enumerate(steps, 1):
    _set_cell(table, r, 0, a, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    _set_cell(table, r, 1, b)
    _set_cell(table, r, 2, c, bold=True, color=BLUE)

tb = _tb(s, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5))
_set_text(tb.text_frame, "最終目標：2026-07-27 POC Demo — 三模組整合展示", size=16, bold=True, color=GRAY_600, alignment=PP_ALIGN.CENTER)

# Big quote
_add_shape_bg(s, Inches(2.0), Inches(6.2), Inches(9.3), Inches(0.8), LIGHT_BLUE_BG, corner_radius=0.05)
tb = _tb(s, Inches(2.2), Inches(6.25), Inches(8.9), Inches(0.7))
_set_text(tb.text_frame, "「好的系統不是取代人，是讓人做得更好。」", size=22, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
# APPENDIX A — Full Feature List
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, NEAR_WHITE)
_section_tag(s, "附錄 A")
_page_number(s, 16)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "完整功能清單（12 項）", size=30, bold=True, color=DARK)

table = _add_table(s, 13, 4, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5.0))
_header_row(table, ["#", "功能", "模組", "優先級"])
features = [
    ("1", "貼文生成引擎", "模組 1", "P0"),
    ("2", "SEO 內容頁", "模組 1", "P2"),
    ("3", "健康問卷系統", "模組 2", "P0"),
    ("4", "問卷結果研判表", "模組 2", "P0"),
    ("5", "產品建議表", "模組 2", "P0"),
    ("6", "名單標籤分級", "模組 2", "P0"),
    ("7", "分流規則引擎", "模組 2", "P1"),
    ("8", "月目標設定", "模組 6", "P0"),
    ("9", "績效看板", "模組 6", "P0"),
    ("10", "培訓執行追蹤", "模組 6", "P1"),
    ("11", "落後提醒", "模組 6", "P1"),
    ("12", "團隊儀表板", "模組 6", "P0"),
]
for r, (a, b, c, d) in enumerate(features, 1):
    bg = WHITE if r % 2 == 1 else NEAR_WHITE
    p_color = BLUE if d == "P0" else (GREEN if d == "P1" else GRAY_400)
    _set_cell(table, r, 0, a, alignment=PP_ALIGN.CENTER, bg=bg)
    _set_cell(table, r, 1, b, bg=bg)
    _set_cell(table, r, 2, c, bg=bg)
    _set_cell(table, r, 3, d, bold=True, color=p_color, alignment=PP_ALIGN.CENTER, bg=bg)

# ══════════════════════════════════════════════════
# APPENDIX B — Tech Stack
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, NEAR_WHITE)
_section_tag(s, "附錄 B")
_page_number(s, 17)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "技術架構", size=30, bold=True, color=DARK)

layers = [
    ("前端", "React 18/19 + Vite + Tailwind v4 + shadcn/ui", BLUE),
    ("後端", "FastAPI + Python 3.12 + uv + Pydantic v2", GREEN),
    ("AI 引擎", "LiteLLM → Gemini 2.5 Flash（可切換）", PURPLE),
    ("資料庫", "PostgreSQL", BLUE),
    ("自動化", "n8n（多平台發布 + 問卷後續觸發）", YELLOW),
    ("部署", "Docker Compose → 雲端（TBD）", GRAY_600),
]
for i, (name, tech, accent) in enumerate(layers):
    y = Inches(2.0 + i * 0.85)
    _add_shape_bg(s, Inches(0.8), y, Inches(2.2), Inches(0.65), accent, corner_radius=0.05)
    tb = _tb(s, Inches(0.9), y + Inches(0.05), Inches(2.0), Inches(0.5))
    _set_text(tb.text_frame, name, size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    tb = _tb(s, Inches(3.3), y + Inches(0.05), Inches(9), Inches(0.5))
    _set_text(tb.text_frame, tech, size=14, color=DARK)

# ══════════════════════════════════════════════════
# APPENDIX C — Success Metrics
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, NEAR_WHITE)
_section_tag(s, "附錄 C")
_page_number(s, 18)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "成功指標（Exit Criteria）", size=30, bold=True, color=DARK)

table = _add_table(s, 10, 4, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.8))
_header_row(table, ["類別", "指標", "目標值", "權重"])
metrics_data = [
    ("模組 1", "AI 貼文可用率", "≥ 80%", "MUST"),
    ("模組 1", "排程發布成功", "≥ 2 平台", "MUST"),
    ("模組 2", "名單分級準確率", "≥ 85%", "MUST"),
    ("模組 2", "問卷→報告全流程", "無需人工", "MUST"),
    ("模組 6", "儀表板即時性", "< 5 分鐘", "MUST"),
    ("技術", "API P95 延遲", "< 500ms", "20%"),
    ("技術", "測試覆蓋率", "≥ 80%", "15%"),
    ("商業", "滿意度評分", "≥ 4.0/5", "15%"),
    ("商業", "UAT 通過率", "≥ 90%", "10%"),
]
for r, (a, b, c, d) in enumerate(metrics_data, 1):
    bg = WHITE if r % 2 == 1 else NEAR_WHITE
    w_color = RED if d == "MUST" else BLUE
    _set_cell(table, r, 0, a, bold=True, bg=bg)
    _set_cell(table, r, 1, b, bg=bg)
    _set_cell(table, r, 2, c, bold=True, color=GREEN, bg=bg)
    _set_cell(table, r, 3, d, bold=True, color=w_color, alignment=PP_ALIGN.CENTER, bg=bg)

# ══════════════════════════════════════════════════
# APPENDIX D — RACI
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, NEAR_WHITE)
_section_tag(s, "附錄 D")
_page_number(s, 19)

tb = _tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "RACI 職責矩陣", size=30, bold=True, color=DARK)

table = _add_table(s, 9, 4, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5))
_header_row(table, ["決策事項", "客戶 (Kuanwei)", "開發團隊", "備註"])
raci_data = [
    ("功能範圍確認", "A（核准）", "R（建議）", "有分歧時客戶拍板"),
    ("技術選型", "I（知情）", "A（負責）", "有疑慮可討論"),
    ("問卷題目定義", "R（負責）", "C（諮詢）", "業務領域知識"),
    ("分流規則設定", "A（核准）", "R（實作）", ""),
    ("UI/UX 方向", "C（諮詢）", "A（負責）", "可提供品牌素材"),
    ("里程碑驗收", "A（核准）", "R（演示）", "Gate Review"),
    ("第三方服務選用", "I（知情）", "A（負責）", "費用實報實銷"),
    ("Go/No-Go", "A（最終決定）", "R（提供建議）", "基於 Exit Criteria"),
]
for r, (a, b, c, d) in enumerate(raci_data, 1):
    bg = WHITE if r % 2 == 1 else NEAR_WHITE
    _set_cell(table, r, 0, a, bold=True, bg=bg)
    _set_cell(table, r, 1, b, color=BLUE, bg=bg, alignment=PP_ALIGN.CENTER)
    _set_cell(table, r, 2, c, color=GREEN, bg=bg, alignment=PP_ALIGN.CENTER)
    _set_cell(table, r, 3, d, size=11, color=GRAY_600, bg=bg)

# ── Save ──
out_path = "/home/os-sunnie.gd.weng/python_workstation/sunny_01/synergy-ai-platform/system_design_docs/poc/Synergy_AI_POC_Pitch.pptx"
prs.save(out_path)
print(f"✅ Saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
